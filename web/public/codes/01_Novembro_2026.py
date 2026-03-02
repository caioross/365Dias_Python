import os
from pptx import Presentation

def extrair_texto_do_pptx(caminho_arquivo):
    """
    Extrai o conteúdo textual de todos os slides de uma apresentação do PowerPoint (.pptx).

    Args:
        caminho_arquivo (str): O caminho para o arquivo .pptx.

    Returns:
        str: Uma string contendo todo o texto extraído dos slides.
    """
    if not os.path.exists(caminho_arquivo):
        raise FileNotFoundError(f"O arquivo {caminho_arquivo} não foi encontrado.")

    texto_extraido = []
    apresentacao = Presentation(caminho_arquivo)

    for slide in apresentacao.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto_extraido.append(shape.text)

    return "\n".join(texto_extraido)

def main():
    caminho_arquivo = 'exemplo.pptx'  # Substitua pelo caminho do seu arquivo .pptx
    try:
        texto = extrair_texto_do_pptx(caminho_arquivo)
        print("Texto extraído do PowerPoint:")
        print(texto)
    except Exception as e:
        print(f"Ocorreu um erro: {e}")

if __name__ == '__main__':
    main()